import os, json, sqlite3, tempfile
from typing import List, Dict

from fastapi import FastAPI, UploadFile, Form
from fastapi.responses import FileResponse

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

import matplotlib.pyplot as plt
from graphviz import Digraph

# ---------------- CONFIG ----------------
DB_PATH = "professor_memory.db"

# ---------------- APP ----------------
app = FastAPI(title="Professor AI Teaching Assistant")

# ---------------- DATABASE ----------------
def init_db():
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS professor_preferences (
            professor_id TEXT PRIMARY KEY,
            layout TEXT,
            diagram_ratio INTEGER,
            math_level TEXT
        )
    """)
    conn.commit()
    conn.close()

def get_preferences(professor_id: str) -> Dict:
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("SELECT layout, diagram_ratio, math_level FROM professor_preferences WHERE professor_id=?",
                (professor_id,))
    row = cur.fetchone()
    conn.close()

    if not row:
        return {
            "layout": "Title and Content",
            "diagram_ratio": 2,
            "math_level": "medium"
        }

    return {
        "layout": row[0],
        "diagram_ratio": row[1],
        "math_level": row[2]
    }

def update_preferences(professor_id, diagram_ratio, math_level):
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("""
        INSERT OR REPLACE INTO professor_preferences
        VALUES (?, ?, ?, ?)
    """, (professor_id, "Title and Content", diagram_ratio, math_level))
    conn.commit()
    conn.close()

# ---------------- PDF ----------------
def extract_text(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# ---------------- TEMPLATE UTIL ----------------
def get_layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[1]

def extract_style(prs: Presentation) -> Dict:
    slide = prs.slides[0]
    run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    return {"font": run.font.name, "size": run.font.size}

def fuse_styles(templates: List[str]) -> Dict:
    fonts, sizes = [], []
    for t in templates:
        prs = Presentation(t)
        style = extract_style(prs)
        fonts.append(style["font"])
        sizes.append(style["size"])
    return {
        "font": max(set(fonts), key=fonts.count),
        "size": max(sizes)
    }

def apply_style(slide, style):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = style["font"]
                r.font.size = style["size"]

# ---------------- DIAGRAM ----------------
def generate_diagram(text: str, path: str):
    dot = Digraph(format="png")
    steps = text.split("->")
    for i in range(len(steps) - 1):
        dot.edge(steps[i].strip(), steps[i + 1].strip())
    dot.render(path, cleanup=True)

# ---------------- MATH ----------------
def render_equation(eq: str, path: str):
    plt.figure()
    plt.text(0.1, 0.5, f"${eq}$", fontsize=22)
    plt.axis("off")
    plt.savefig(path, bbox_inches="tight")
    plt.close()

# ---------------- LLM SLIDE PLANNER (STUB) ----------------
def llm_plan_slides(notes: str, prefs: Dict) -> List[Dict]:
    # Replace this with Gemini / OpenAI call
    chunks = [c for c in notes.split("\n\n") if len(c) > 60][:6]
    slides = []

    for i, c in enumerate(chunks):
        slides.append({
            "title": f"Concept {i+1}",
            "content": c[:400],
            "diagram": "Input -> Process -> Output" if i % prefs["diagram_ratio"] == 0 else None,
            "equation": r"\hat{y} = Wx + b" if prefs["math_level"] != "low" else None
        })

    return slides

# ---------------- PPT GENERATION ----------------
def create_ppt(templates: List[str], slides: List[Dict], prefs: Dict, out: str):
    prs = Presentation(templates[0])
    style = fuse_styles(templates)
    layout = get_layout_by_name(prs, prefs["layout"])

    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s["title"]
        slide.placeholders[1].text = s["content"]

        if s["diagram"]:
            d_path = f"diagram_{i}"
            generate_diagram(s["diagram"], d_path)
            slide.shapes.add_picture(d_path + ".png", Inches(4.5), Inches(1.5), width=Inches(4))

        if s["equation"]:
            m_path = f"math_{i}.png"
            render_equation(s["equation"], m_path)
            slide.shapes.add_picture(m_path, Inches(1), Inches(4.5), width=Inches(3))

        apply_style(slide, style)

    prs.save(out)

# ---------------- API ----------------
@app.post("/generate_ppt")
async def generate_ppt(
    pdf: UploadFile,
    templates: List[UploadFile],
    professor_id: str = Form(...)
):
    init_db()

    with tempfile.NamedTemporaryFile(delete=False) as p:
        p.write(await pdf.read())
        pdf_path = p.name

    template_paths = []
    for t in templates:
        tmp = tempfile.NamedTemporaryFile(delete=False)
        tmp.write(await t.read())
        template_paths.append(tmp.name)

    notes = extract_text(pdf_path)
    prefs = get_preferences(professor_id)
    slides = llm_plan_slides(notes, prefs)

    output = f"{professor_id}_lecture.pptx"
    create_ppt(template_paths, slides, prefs, output)

    return FileResponse(output, filename=output)

@app.post("/feedback")
def feedback(
    professor_id: str,
    more_diagrams: bool,
    less_math: bool
):
    update_preferences(
        professor_id,
        diagram_ratio=1 if more_diagrams else 3,
        math_level="low" if less_math else "high"
    )
    return {"status": "preferences updated"}

# ---------------- RUN ----------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
