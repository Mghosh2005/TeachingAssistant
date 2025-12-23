"""
AI Teaching Assistant – Full Deployable Single-File System
---------------------------------------------------------
Features:
1. LLM-based slide generation (Gemini / OpenAI compatible)
2. FAISS RAG to prevent hallucination
3. FastAPI backend
4. PPT template-based rendering

Run:
pip install fastapi uvicorn python-pptx langchain faiss-cpu pymupdf google-generativeai python-dotenv
uvicorn app:app --reload
"""

import os
import fitz
from fastapi import FastAPI, UploadFile, File
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.embeddings import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from dotenv import load_dotenv

# ------------------------
# ENV SETUP
# ------------------------
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# ------------------------
# FASTAPI APP
# ------------------------
app = FastAPI(title="AI Teaching Assistant")

# ------------------------
# PDF TEXT EXTRACTION
# ------------------------
def extract_text_from_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return "".join(page.get_text() for page in doc)

# ------------------------
# RAG SETUP
# ------------------------
def build_vector_store(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    chunks = splitter.split_text(text)

    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    return FAISS.from_texts(chunks, embeddings)

# ------------------------
# SLIDE PLANNER (LLM)
# ------------------------
def plan_slides(model, context):
    prompt = f"""
You are a university professor.
Create a teaching slide plan.

Rules:
- One concept per slide
- Teaching friendly
- Exam oriented

Context:
{context[:1500]}

Output JSON list:
[{{"title": "", "type": ""}}]
"""
    response = model.generate_content(prompt)
    return eval(response.text)

# ------------------------
# SLIDE CONTENT GENERATOR (LLM)
# ------------------------
def generate_slide(model, title, context):
    prompt = f"""
Create slide content.

Title: {title}
Rules:
- Max 5 bullets
- Academic language
- Clear explanation

Context:
{context}

Output JSON:
{{"title": "", "bullets": []}}
"""
    response = model.generate_content(prompt)
    return eval(response.text)

# ------------------------
# PPT RENDERER (TEMPLATE-BASED)
# ------------------------
def render_ppt(slides, template="template.pptx"):
    prs = Presentation(template if os.path.exists(template) else None)

    for slide in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide["title"]
        s.placeholders[1].text = "\n".join(slide["bullets"])

    output = "lecture.pptx"
    prs.save(output)
    return output

# ------------------------
# API ENDPOINT
# ------------------------
@app.post("/generate-ppt/")
async def generate_ppt(file: UploadFile = File(...)):
    pdf_bytes = await file.read()

    # 1. Extract text
    text = extract_text_from_pdf(pdf_bytes)

    # 2. RAG
    vector_store = build_vector_store(text)
    context_docs = vector_store.similarity_search("entire topic", k=5)
    context = " ".join(doc.page_content for doc in context_docs)

    # 3. LLM
    model = genai.GenerativeModel("gemini-pro")

    slide_plan = plan_slides(model, context)

    slides = []
    for slide_meta in slide_plan:
        slide = generate_slide(model, slide_meta["title"], context)
        slides.append(slide)

    # 4. PPT
    ppt_path = render_ppt(slides)

    return {
        "status": "success",
        "ppt_file": ppt_path
    }

