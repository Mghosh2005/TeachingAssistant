import os, json, sqlite3, tempfile
from typing import List, Dict
from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI, UploadFile, Form
from fastapi.responses import FileResponse
import re
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import requests
from bs4 import BeautifulSoup
import matplotlib.pyplot as plt
from graphviz import Digraph
from fastapi import HTTPException
import shutil
from pptx.util import Pt

def split_into_sentences(text: str) -> List[str]:
    # Normalize whitespace so pdf line breaks don't ruin sentence splitting
    text = re.sub(r'\s+', ' ', text).strip()
    # Simple sentence split (works well for English lecture notes)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    # If too short or no punctuation, split on commas/semicolons as fallback
    if len(sentences) == 1:
        sentences = re.split(r'[,\n;]\s*', text)
    # Filter tiny fragments
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    return sentences

def extract_text_from_url(url: str) -> str:
    headers = {
        "User-Agent": "Mozilla/5.0"
    }
    resp = requests.get(url, headers=headers, timeout=10)
    resp.raise_for_status()

    soup = BeautifulSoup(resp.text, "html.parser")

    # Remove scripts & styles
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    text = soup.get_text(separator=" ")
    text = re.sub(r'\s+', ' ', text).strip()

    return text


# ---------------- CONFIG ----------------
DB_PATH = "professor_memory.db"

# ---------------- APP ----------------
app = FastAPI(title="Professor AI Teaching Assistant")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------- DATABASE ----------------
def init_db():
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS professor_preferences (
            professor_id TEXT PRIMARY KEY,
            layout TEXT,
            diagram_ratio INTEGER,
            math_level TEXT
        )
    """)
    conn.commit()
    conn.close()

def get_preferences(professor_id: str) -> Dict:
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("SELECT layout, diagram_ratio, math_level FROM professor_preferences WHERE professor_id=?",
                (professor_id,))
    row = cur.fetchone()
    conn.close()

    if not row:
        return {
            "layout": "Title and Content",
            "diagram_ratio": 2,
            "math_level": "medium"
        }

    return {
        "layout": row[0],
        "diagram_ratio": row[1],
        "math_level": row[2]
    }

def update_preferences(professor_id, diagram_ratio, math_level):
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("""
        INSERT OR REPLACE INTO professor_preferences
        VALUES (?, ?, ?, ?)
    """, (professor_id, "Title and Content", diagram_ratio, math_level))
    conn.commit()
    conn.close()

# ---------------- PDF ----------------
def extract_text(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# ---------------- TEMPLATE UTIL ----------------
def get_layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[1]



def extract_style(prs: Presentation) -> Dict:
    """
    Safely extract font name and size from any available text run
    in the presentation. Falls back to defaults if not found.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            for para in tf.paragraphs:
                if para.runs:
                    run = para.runs[0]
                    return {
                        "font": run.font.name or "Calibri",
                        "size": run.font.size or Pt(24)
                    }

    # Absolute fallback (never crash)
    return {
        "font": "Calibri",
        "size": Pt(24)
    }


def fuse_styles(templates: List[str]) -> Dict:
    fonts, sizes = [], []
    for t in templates:
        prs = Presentation(t)
        style = extract_style(prs)
        fonts.append(style["font"])
        sizes.append(style["size"])
    return {
        "font": max(set(fonts), key=fonts.count),
        "size": max(sizes)
    }

def apply_style(slide, style):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = style["font"]
                r.font.size = style["size"]

# ---------------- DIAGRAM ----------------
def generate_diagram(text: str, path: str):
    dot = Digraph(format="png")
    steps = text.split("->")
    for i in range(len(steps) - 1):
        dot.edge(steps[i].strip(), steps[i + 1].strip())
    dot.render(path, cleanup=True)

# ---------------- MATH ----------------
def render_equation(eq: str, path: str):
    plt.figure()
    plt.text(0.1, 0.5, f"${eq}$", fontsize=22)
    plt.axis("off")
    plt.savefig(path, bbox_inches="tight")
    plt.close()

# ---------------- LLM SLIDE PLANNER (STUB) ----------------
def llm_plan_slides(notes: str, prefs: Dict, num_slides: int = 6) -> List[Dict]:
    """
    Lightweight deterministic planner:
    - split notes into sentences
    - distribute sentences roughly evenly across num_slides
    - create slide dicts with title (first sentence / short headline) and content (joined)
    """
    sentences = split_into_sentences(notes)
    n = len(sentences)
    if n == 0:
        return [{"title": "Empty", "content": "No content found.", "diagram": None, "equation": None}]

    # Don't create more slides than we have sentences (optional)
    num_slides = min(num_slides, max(1, n))

    base = n // num_slides
    extra = n % num_slides

    slides = []
    start = 0
    for i in range(num_slides):
        count = base + (1 if i < extra else 0)
        chunk_sents = sentences[start:start+count]
        start += count

        title = chunk_sents[0]
        # Make a short headline for the title (first 6 words)
        title_short = " ".join(title.split()[:6]).rstrip(' ,.;:') 
        content = "\n".join(chunk_sents)

        slides.append({
            "title": title_short,
            "content": content,
            # diagram/equation heuristics: use teacher prefs for frequency
            "diagram": "Input -> Process -> Output" if (i % max(1, prefs.get("diagram_ratio",2)) == 0) else None,
            "equation": r"\hat{y} = Wx + b" if prefs.get("math_level","medium") != "low" and "y =" in notes else None
        })

    return slides

@app.post("/api/generate-outline/")
async def generate_outline(
    pdf: UploadFile = None,
    professor_id: str = Form(...),
    num_slides: int = Form(10),
    urls: str = Form("")
):
    init_db()

    notes = ""

    # ---------- PDF INPUT ----------
    if pdf is not None:
        with tempfile.NamedTemporaryFile(delete=False) as p:
            p.write(await pdf.read())
            pdf_path = p.name

        notes += extract_text(pdf_path)

    # ---------- URL INPUT ----------
    if urls:
        for url in urls.splitlines():
            url = url.strip()
            if not url:
                continue

            try:
                notes += "\n" + extract_text_from_url(url)
            except Exception as e:
                print(f"[URL ERROR] {url}: {e}")

    # ---------- SAFETY CHECK ----------
    if not notes.strip():
        return {
            "outline": [{
                "title": "No Content Provided",
                "points": [
                    "Please upload a PDF or provide one or more valid URLs."
                ]
            }]
        }

    # ---------- SLIDE PLANNING ----------
    prefs = get_preferences(professor_id)

    slides = llm_plan_slides(
        notes=notes,
        prefs=prefs,
        num_slides=num_slides
    )

    # ---------- OUTLINE RESPONSE ----------
    outline = [
        {
            "title": slide["title"],
            "points": slide["content"].split("\n")[:5]
        }
        for slide in slides
    ]

    return {"outline": outline}


# ---------------- PPT GENERATION ----------------
def create_ppt(templates: List[str], slides: List[Dict], prefs: Dict, out: str):
    prs = Presentation(templates[0])
    style = fuse_styles(templates)
    layout = get_layout_by_name(prs, prefs["layout"])

    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(layout)

        # ----- TITLE -----
        if slide.shapes.title:
            slide.shapes.title.text = s["title"]

        # ----- CONTENT (SAFE) -----
        content_written = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.clear()
                shape.text_frame.text = s["content"]
                content_written = True
                break

        if not content_written:
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(4)
            )
            txBox.text_frame.text = s["content"]

        # ----- DIAGRAM -----
        if s.get("diagram"):
            d_path = f"diagram_{i}"
            generate_diagram(s["diagram"], d_path)
            slide.shapes.add_picture(
                d_path + ".png",
                Inches(4.5),
                Inches(1.5),
                width=Inches(4)
            )

        # ----- EQUATION -----
        if s.get("equation"):
            m_path = f"math_{i}.png"
            render_equation(s["equation"], m_path)
            slide.shapes.add_picture(
                m_path,
                Inches(1),
                Inches(4.5),
                width=Inches(3)
            )

        # ----- STYLE -----
        apply_style(slide, style)

    prs.save(out)

# ---------------- API ----------------
@app.post("/api/generate-slides/")
async def generate_slides(data: Dict):
    """
    Receives edited outline from UI
    Returns full slide content for professor review
    """
    slides = []

    for s in data["outline"]:
        slides.append({
            "title": s["title"],
            "content": "\n".join(s["points"]),
            "approved": True
        })

    return {"slides": slides}

@app.post("/api/render-ppt/")
async def render_ppt(
    slides: str = Form(...),
    templates: List[UploadFile] = None,
    professor_id: str = Form(...)
):
    init_db()

    # ---------- Templates ----------
    template_paths: List[str] = []
    if templates:
        for t in templates:
            # Basic server-side extension check
            if not t.filename.lower().endswith(".pptx"):
                # skip non-pptx or raise
                continue
            # Create a temp file and save uploaded template
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            try:
                # read chunked to avoid big memory use
                with open(tmp.name, "wb") as out_file:
                    shutil.copyfileobj(t.file, out_file)
                template_paths.append(tmp.name)
            finally:
                # t.file is an SpooledTemporaryFile created by Starlette/FastAPI;
                # do not close t.file too early if reusing, but it's ok to let it be GC'd
                pass

    # Fallback to built-in default template if no templates uploaded
    if not template_paths:
        # Use the project default_template.pptx file (make sure path correct)
        default_path = os.path.join(os.path.dirname(__file__), "default_template.pptx")
        if not os.path.exists(default_path):
            raise HTTPException(status_code=500, detail="Default template not found on server.")
        template_paths = [default_path]

    # ---------- Slides ----------
    prefs = get_preferences(professor_id)

    try:
        slides_list = json.loads(slides)
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid slides JSON")

    # ---------- Generate PPT ----------
    output = f"{professor_id}_lecture.pptx"
    create_ppt(template_paths, slides_list, prefs, output)

    # Optionally clean up uploaded template temp files here (if you want to remove them)
    # for p in template_paths:
    #     if p != default_path:
    #         try: os.remove(p)
    #         except: pass

    return FileResponse(output, filename=output)

@app.post("/feedback")
def feedback(
    professor_id: str,
    more_diagrams: bool,
    less_math: bool
):
    update_preferences(
        professor_id,
        diagram_ratio=1 if more_diagrams else 3,
        math_level="low" if less_math else "high"
    )
    return {"status": "preferences updated"}

# ---------------- RUN ----------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)